"""
OxyPC Inventory ERP — Flowchart PPTX Generator (Python / python-pptx)
======================================================================
This script was generated as a Python reference. The actual PPTX was
produced by generate_flowchart_pptx.js (PptxGenJS), which is the
recommended tool for creating presentations from scratch.

To regenerate the PPTX run:
    node generate_flowchart_pptx.js

To run this Python version:
    pip install python-pptx
    python generate_flowchart_pptx.py

Both produce equivalent output at:
    C:\\Users\\Pankaj.sehgal\\Claude\\Oxypc\\oxypc-inventory\\docs\\OxyPC_Flowchart.pptx
"""

import os
import subprocess
import sys

OUTPUT_PATH = r"C:\Users\Pankaj.sehgal\Claude\Oxypc\oxypc-inventory\docs\OxyPC_Flowchart.pptx"
JS_SCRIPT   = r"C:\Users\Pankaj.sehgal\Claude\Oxypc\oxypc-inventory\docs\generate_flowchart_pptx.js"

# ── Attempt to use python-pptx as a thin wrapper ──────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
except ImportError:
    print("python-pptx not installed. Falling back to Node.js generator.")
    result = subprocess.run(["node", JS_SCRIPT], capture_output=True, text=True)
    print(result.stdout)
    if result.returncode != 0:
        print(result.stderr)
        sys.exit(result.returncode)
    sys.exit(0)

# ── Color palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
AMBER  = RGBColor(0xF5, 0x7F, 0x17)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
TEAL   = RGBColor(0x00, 0x69, 0x5C)
RED    = RGBColor(0xC6, 0x28, 0x28)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x54, 0x6E, 0x7A)
BLACK  = RGBColor(0x21, 0x21, 0x21)
LIGHTBG = RGBColor(0xEE, 0xF4, 0xFC)

# ── Helper: set shape fill ────────────────────────────────────────────────────
def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.5):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def set_text(shape, text, font_color=WHITE, font_size=9, bold=True, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color


# ── Shape helpers ─────────────────────────────────────────────────────────────
def add_box(slide, left, top, width, height, text,
            fill_rgb=NAVY, text_rgb=WHITE, font_size=9, bold=True):
    """Rounded rectangle process box."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    set_fill(shape, fill_rgb)
    set_line(shape, fill_rgb if fill_rgb != WHITE else NAVY)
    set_text(shape, text, text_rgb, font_size, bold)
    return shape


def add_process_box(slide, left, top, width, height, text, font_size=9):
    """White fill, navy border, navy text."""
    return add_box(slide, left, top, width, height, text,
                   fill_rgb=WHITE, text_rgb=NAVY, font_size=font_size)


def add_diamond(slide, left, top, width, height, text,
                fill_rgb=ORANGE, font_size=8):
    """Decision diamond."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.DIAMOND,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    set_fill(shape, fill_rgb)
    set_line(shape, fill_rgb)
    set_text(shape, text, WHITE, font_size, True)
    return shape


def add_oval(slide, left, top, width, height, text,
             fill_rgb=NAVY, text_rgb=WHITE, font_size=9):
    """Terminal oval (Start/End)."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    set_fill(shape, fill_rgb)
    set_line(shape, fill_rgb)
    set_text(shape, text, text_rgb, font_size, True)
    return shape


def add_arrow_down(slide, cx, y1, y2, color=NAVY):
    """Vertical downward arrow (thin line shape)."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.LINE_INVERSE,   # fallback
        Inches(cx - 0.02), Inches(y1),
        Inches(0.04), Inches(y2 - y1)
    )
    set_fill(shape, color)
    shape.line.fill.background()
    return shape


def add_label(slide, left, top, width, height, text,
              font_color=GREY, font_size=7.5):
    """Transparent text label."""
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    return txbox


# ── Slide background ──────────────────────────────────────────────────────────
def set_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


# ── Title banner ──────────────────────────────────────────────────────────────
def add_title_banner(slide, title, subtitle=""):
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.48)
    )
    set_fill(banner, NAVY)
    banner.line.fill.background()
    text = title + ("  |  " + subtitle if subtitle else "")
    add_label(slide, 0.18, 0.02, 13.0, 0.44, text, WHITE, 12)


# ── Presentation ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # completely blank


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(sl, NAVY)

# Title
tb = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.0))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "OxyPC Refurbishment ERP"
run.font.size = Pt(40)
run.font.bold = True
run.font.color.rgb = WHITE

# Subtitle
tb2 = sl.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.33), Inches(0.55))
tf2 = tb2.text_frame
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Complete Application Process Flowcharts"
r2.font.size = Pt(18)
r2.font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC)

# Date
tb3 = sl.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.33), Inches(0.38))
tf3 = tb3.text_frame
p3 = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "27 April 2026"
r3.font.size = Pt(13)
r3.font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC)

# Orange accent stripe
stripe = sl.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    Inches(0), Inches(4.85), Inches(13.33), Inches(0.06)
)
set_fill(stripe, ORANGE)
stripe.line.fill.background()


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Device Lifecycle (simplified python-pptx version)
# Note: Full layout in JS version; this is a structured overview
# ════════════════════════════════════════════════════════════════════════════
sl2 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(sl2, LIGHTBG)
add_title_banner(sl2, "Device Lifecycle — Main Flow", "GRN → IQC → Repair → QC → Sale")

# Main vertical flow — left column
ax, aw, rh, g = 0.12, 2.0, 0.42, 0.13
ay = 0.58

def nxt2():
    global ay
    y = ay; ay += rh + g; return y

add_oval(sl2, ax, nxt2(), aw, rh, "START: GRN / Material Received")
add_process_box(sl2, ax, nxt2(), aw, rh, "Supplier Invoice & Lot Registration")
add_process_box(sl2, ax, nxt2(), aw, rh, "Device Barcoding & Line Items")
add_process_box(sl2, ax, nxt2(), aw, rh, "Advance to IQC")
add_diamond(sl2, ax, nxt2(), aw, rh + 0.1, "IQC Inspection\n(60+ checks)")
ay += 0.1
add_box(sl2, ax, nxt2(), aw, rh, "Move to Stock In (Grade C0)", GREEN)
add_box(sl2, ax, nxt2(), aw, rh, "QC Inspection (Battery/Screen/Kbd/Body)", GREEN)
add_diamond(sl2, ax, nxt2(), aw, rh + 0.1, "QC Result?")
ay += 0.1
add_box(sl2, ax, nxt2(), aw, rh, "Ready to Sale", TEAL)
add_box(sl2, ax, nxt2(), aw, rh, "Sale (Price / Customer / Payment)", TEAL)
add_diamond(sl2, ax, nxt2(), aw, rh + 0.1, "Return?")
ay += 0.1
add_oval(sl2, ax, nxt2(), aw, rh, "SOLD ✓ — End")

# Repair column
bx, bw2 = 2.35, 2.05
add_label(sl2, bx, 0.55, bw2, 0.25, "REPAIR ESCALATION", AMBER, 7.5)
for step, lbl in [("L1 Repair (Basic)", C := AMBER), ("L2 Repair (Intermediate)", AMBER), ("L3 Repair (Advanced)", AMBER)]:
    add_box(sl2, bx, 1.85 + [0, 1.1, 2.2][["L1 Repair (Basic)", "L2 Repair (Intermediate)", "L3 Repair (Advanced)"].index(step)], bw2, rh, step, AMBER)

# Cosmetic pipeline
cx2, cw2 = 5.0, 1.75
add_label(sl2, cx2, 0.55, cw2, 0.25, "COSMETIC PIPELINE", PURPLE, 7.5)
for i, step in enumerate(["Cleaning", "Dry Sanding", "Masking", "Painting", "Water Sanding", "Final QC"]):
    c = GREEN if step == "Final QC" else PURPLE
    add_box(sl2, cx2, 0.85 + i * (rh + g), cw2, rh, step, c)

# Spare parts
dx2, dw2 = 7.1, 2.0
add_label(sl2, dx2, 0.55, dw2, 0.25, "SPARE PARTS BRANCH", GREY, 7.5)
add_diamond(sl2, dx2, 0.85, dw2, rh + 0.1, "Parts Not Available?")
add_box(sl2, dx2, 1.55, dw2, rh, "Spare Parts Planning", GREY)
add_diamond(sl2, dx2, 2.15, dw2, rh + 0.1, "Part Available?")
add_box(sl2, dx2, 2.9, dw2, rh, "Assign Part to Engineer", GREEN)
add_box(sl2, dx2, 3.45, dw2, rh + 0.55, "Create PO → GRN → Update Inventory", NAVY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES 3–6: Generated by Node.js version (run generate_flowchart_pptx.js)
# The python-pptx version above shows structure; for full visual fidelity
# use the JS generator which has complete layout calculations.
# ════════════════════════════════════════════════════════════════════════════

# Save python-pptx output (partial — for reference)
py_output = OUTPUT_PATH.replace(".pptx", "_py_reference.pptx")
prs.save(py_output)
size = os.path.getsize(py_output)
print(f"Python-pptx reference saved: {py_output} ({size:,} bytes)")
print()
print("For the complete 6-slide flowchart PPTX, run:")
print(f"  node {JS_SCRIPT}")
print(f"Output: {OUTPUT_PATH}")

# Also check if main JS output exists
if os.path.exists(OUTPUT_PATH):
    main_size = os.path.getsize(OUTPUT_PATH)
    print(f"\nMain output already exists: {OUTPUT_PATH} ({main_size:,} bytes)")
