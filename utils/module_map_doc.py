"""Module & Interconnectivity Map — renders the system reference as PDF and PPTX.

The content below is the same map published as a system reference on 7 Aug 2026:
module names come from PERM_MODULES in routers/master.py, and the flows are derived
from the foreign-key graph across the model files plus the stage machine in
services/control_engine.py. Keep it in sync when modules or flows change.
"""
from __future__ import annotations

from io import BytesIO

from reportlab.lib import colors
from reportlab.lib.pagesizes import LETTER, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    KeepTogether, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
)

_NAVY = colors.HexColor("#1F3864")
_ACCENT = colors.HexColor("#2E75B6")
_LIGHT = colors.HexColor("#EEF3F8")
_MUTED = colors.HexColor("#5A6472")
_GREEN = colors.HexColor("#2E6B4C")
_RED = colors.HexColor("#9E3A2C")

SOURCE_REF = "oxypc-inventory"
GENERATED_LABEL = "88 modules · 74 routers · 119 tables · 142 FK edges · 14 services · 20 stages"

HEADLINE_STATS = [
    ("88", "Modules"), ("74", "Routers"), ("119", "Tables"),
    ("142", "FK edges"), ("14", "Services"), ("20", "Stages"),
]

# ── Content ───────────────────────────────────────────────────────────────────

SPINE_INTRO = (
    "devices.current_stage is the axis this system turns on. Transitions are not "
    "hard-coded: they are rows in allowed_transitions, validated by the control engine on "
    "every move and written to stage_movements as an append-only history. Change the table, "
    "change the process — no deploy required."
)

SPINE_STAGES = [
    "grn → iqc → stock_in → trc_production",
    "trc_production → l1 → l2 → l3 (escalation order enforced server-side)",
    "l1 / l2 / l3 → qc_check (stress test)",
    "qc_check → cleaning → putty → dry_sanding → masking → painting → water_sanding",
    "water_sanding → final_qc → ready_to_sale → sold → returned",
    "l3 or final_qc → scrapped → scrap_for_sale",
]

SPINE_NOTE = (
    "Six cosmetic stages, not five. Cleaning → Putty → Dry Sanding → Masking → Painting → "
    "Water Sanding is a finer-grained paint line than most refurbishment systems model. It is "
    "the reason per-stage aging and labour costing can be attributed accurately rather than "
    "lumped into one \"refurb\" bucket."
)

MODULE_GROUPS = [
    ("Procure & Intake", ["crm_sourcing", "crm_purchase_orders", "grn", "grn_records",
                          "grn_post_iqc", "lots", "iqc"]),
    ("Stock & Movement", ["stock", "devices", "entity_movement", "transfers",
                          "move_device_internal", "inventory_requests", "model_requested"]),
    ("Production & Repair", ["production_manager", "repair_l1", "repair_l2", "repair_l3",
                             "repair_l3l4", "qc_check", "cosmetic", "cosmetic_finalqc",
                             "workid_status", "scrap_products"]),
    ("Spare Parts", ["parts_dashboard", "spare_parts", "spare_parts_purchase",
                     "parts_sale_request", "parts_tracking", "parts_consumption"]),
    ("Sales & Fulfilment", ["sales", "sales_list", "dispatch", "gate_pass", "returns",
                            "quotations", "ready_to_sale_parts", "part_sales"]),
    ("Dealers & Trade Partner", ["dealers", "trade_partner", "partner_payments",
                                 "assign_dealer_leads", "model_requests"]),
    ("CRM & Demand", ["crm_dashboard", "crm_contacts", "crm_sales_opp", "crm_quotes",
                      "crm_price_matrix", "crm_analytics", "crm_assign_leads", "telecalling",
                      "telesales_dashboard", "whatsapp", "market", "care_support"]),
    ("Warehouse & Location", ["locations", "location_master", "location_gaps",
                              "location_audit", "location_trash"]),
    ("Finance", ["finance", "finance_supplier", "finance_customer"]),
    ("Reporting", ["reports", "report_sales", "report_bizpl", "report_stage", "report_aging",
                   "report_overdue", "report_receivables"]),
    ("People", ["attendance", "attendance_report", "attendance_config"]),
    ("Platform & Governance", ["dashboard", "stage_control", "aging_tracker", "stage_audit_log",
                               "system_audit_log", "wa_audit_log", "sidebar_config",
                               "landing_pages", "admin_users", "admin_master",
                               "company_settings", "cost_config", "terms_conditions",
                               "qa", "manuals"]),
]

# (code, title, summary, [flow steps], key-joints note, warning or None)
FLOWS = [
    ("F1", "Procurement → intake → inspection",
     "A sourcing deal becomes a lot; the lot receives goods against an invoice; devices are "
     "born at GRN and inspected into IQC. This is the only flow that creates device records "
     "in bulk.",
     ["crm_contacts (supplier) → crm_sourcing_deals",
      "crm_sourcing_deals → crm_purchase_orders and → lots",
      "crm_purchase_orders → grn / parts_grn → lot_line_items",
      "lots + lot_line_items → devices → iqc_inspections",
      "bulk_upload (CSV) → devices"],
     "crm_sourcing_deals → lots is a real FK, so a lot's commercial origin is traceable. "
     "devices → lots and devices → lot_line_items anchor every unit to its purchase. Bulk "
     "upload writes straight into devices and iqc_inspections.",
     None),
    ("F2", "Production & repair loop",
     "The longest flow, and the only one that can send a unit backwards. Work orders carry "
     "the assignment; repair jobs and attempts carry the history; the stress test and final "
     "QC are the two gates that decide whether a unit reaches sale or scrap.",
     ["devices → buckets → production_manager → work_orders",
      "work_orders → repair_jobs → repair_attempts and → part_requests",
      "repair_jobs → qc_checks (stress) → cosmetic (6 stages) → final_qc",
      "final_qc pass → ready_to_sale; final_qc fail → scrapped",
      "stage_movements records every hop against the device"],
     "work_orders → devices, stock_transfers, users — the work order is what ties a person to "
     "a unit to a movement. Escalation order (L1 before L2 before L3) is enforced by the "
     "control engine, not by the UI.",
     None),
    ("F3", "Spare parts — purchase, consume, resell",
     "Parts have their own intake and their own sales channel, and they meet the device flow "
     "at exactly one point: consumption during repair. That single join is what makes parts "
     "cost land in device cost.",
     ["parts_grn and spare_parts_purchases → spare_parts",
      "spare_parts → part_requests → part_sourcing_requests",
      "part_requests → spare_parts_consumption → devices + repair_jobs",
      "spare_parts → part_sale_requests → part_sales",
      "spare_parts → spare_parts_ledger; ram_tracking → devices"],
     "spare_parts_consumption → devices, lots, repair_jobs, spare_parts is the four-way join "
     "that feeds Parts Consumed on the Business P&L. ram_tracking follows RAM sticks "
     "individually, separate from generic part stock.",
     None),
    ("F4", "Demand — CRM, telecalling, quotes",
     "Demand enters through four doors — contacts, social lead assignment, telecalling and "
     "WhatsApp — and converges on quotes and opportunities. crm_contacts is the second-biggest "
     "hub in the system.",
     ["crm_lead_groups → crm_leads → crm_lead_calls",
      "crm_leads → crm_contacts → crm_activities, crm_contact_numbers",
      "crm_contacts → crm_sales_opportunities and → crm_quotes → crm_quote_items",
      "telecalling_records → crm_contacts and → crm_quotes",
      "whatsapp_messages → crm_quotes; crm_contacts → dealers",
      "crm_grade_price_matrix supplies pricing to quotes (soft link)"],
     "dealers → crm_contacts means a dealer is a specialisation of a contact, not a parallel "
     "record. telecalling_records reaches into lots, lot line items, quotes, dealers and "
     "orders — six FKs, the widest single table in the CRM area.",
     None),
    ("F5", "Sale, dispatch and return",
     "Sale is per-device. Everything downstream — warranty, care, returns, receivables — hangs "
     "off the sale row, which is why a device that ships without a sale record silently "
     "disappears from four other modules.",
     ["ready_to_sale → sales → dispatch → gate_pass",
      "sales → returns",
      "sales → care_warranties and → care_device_pairings → care_support_tickets",
      "sales → customer_receipts",
      "telecaller_dispatch_requests → devices → sales"],
     "sales → devices is the anchor for warranty, care pairing, returns and receivables.",
     "Open defect on this path — Backlog #157: New Sale does not always move a device from "
     "Ready-to-Sale to Sold. A stuck stage here shows up as inventory that looks available but "
     "is already committed."),
    ("F6", "Trade Partner — the B2B channel",
     "A self-contained dealer-facing portal on top of the same inventory. Lots are published "
     "as listings; dealers view, bid, book and pay. Twelve tables that touch core inventory at "
     "only two points.",
     ["lots → partner_listings; devices → partner_listing_devices",
      "dealers → partner_listing_views",
      "dealers → partner_bids → partner_bid_documents, partner_bid_payments",
      "dealers → partner_bookings → partner_payment_proofs",
      "lots → lot_dealer_visibility, lot_booking_requests",
      "partner_floor_config sets the floor price applied to bids"],
     "Only partner_listings → lots and partner_listing_devices → devices cross into core "
     "inventory. That narrow waist is why the portal could ship as an MVP without "
     "destabilising the rest of the system.",
     None),
    ("F7", "Warehouse — location, movement, audit",
     "Physical position is tracked separately from process stage. A device has a stage and a "
     "location, and the audit flow reconciles what the system believes against what is on the "
     "rack.",
     ["storage_locations → devices and → buckets",
      "devices + storage_locations → device_location_logs",
      "buckets + devices + storage_locations → stock_transfers",
      "inventory_audits → audit_scan_items → devices, storage_locations",
      "audit_scan_items → gap alerts (found / missing / extra)"],
     "audit_scan_items carries four FKs — device, audit, location and user — which is what "
     "allows a discrepancy to be attributed to a person and a place, not just counted.",
     None),
    ("F8", "Cost and money roll-up",
     "Every cost component is captured where it is incurred, then aggregated. This is the flow "
     "where a missing component silently understates COGS rather than throwing an error.",
     ["lots (buying price) → device_costing",
      "spare_parts_consumption → device_costing",
      "labour / work_orders → device_costing",
      "device_costing + sales → report_bizpl (Business P&L)",
      "lots → Lot P&L report",
      "supplier_payments and customer_receipts → finance",
      "cost_config supplies rates; business_pl_overrides applies manual adjustments"],
     "cost_config → device_costing and business_pl_overrides → Business P&L are the two "
     "configuration inputs that change reported margin without changing any transaction.",
     "Known gap on this path — Risk 1 from the April audit (P&L COGS completeness) is still "
     "open. Labour was added and Parts Consumed now has its own card, but any cost component "
     "not wired into device_costing inflates margin without warning."),
]

SERVICES = [
    ("control_engine", "Validates stage transitions, blocks sale of non-ready devices, "
                       "enforces L1→L2→L3 escalation order", "Every stage-changing module"),
    ("audit_engine", "Writes audit_logs with actor, table, record, before/after", "All write paths"),
    ("cost_engine", "Builds device_costing from lot, parts and labour", "F8"),
    ("parts_required", "Maps a repair issue to the parts it needs", "F2, F3"),
    ("aging_tracker", "Maintains device_aging per stage", "Reporting"),
    ("event_bus + webhook_dispatcher", "Publishes domain events to registered webhooks",
     "External integrations"),
    ("notifications", "In-app notification fan-out", "All modules"),
    ("care_service", "Post-sale warranty, pairing, ticket lifecycle", "F5"),
    ("partner_service", "Listing, bid and booking rules", "F6"),
    ("opportunity_lot", "Matches CRM demand to available lots", "F4 ↔ F1"),
    ("invoice_parser · po_pdf", "Document ingest and generation", "F1"),
    ("call_service", "Telecalling session and outcome handling", "F4"),
]

GOVERNANCE_NOTE = (
    "Governance sits alongside: master_data supplies every dropdown; role_module_permissions, "
    "custom_roles and role_additional_permissions gate module access; app_settings holds "
    "sidebar labels and page titles; stage_master and allowed_transitions define the process "
    "itself."
)

COUPLING = [
    ("devices", "25", "The system's centre of gravity. Any column change is a 25-table blast radius."),
    ("dealers", "17", "Channel spine — orders, calls, credit notes, bids, bookings, WhatsApp groups."),
    ("users", "12", "Attribution for audit, attendance, QC, repair, location logs."),
    ("crm_contacts", "11", "Serves as supplier and customer; dealers specialise from it."),
    ("lots", "11", "Commercial unit of purchase — costing, visibility, bookings, payments."),
    ("spare_parts", "7", "Both a consumable and a sellable good, so it spans two flows."),
]

SOFT_LINKS = (
    "devices.entity is a string validated against master data, not an FK — which is why a typo "
    "would corrupt Entity Movement counts silently, and why bulk upload now rejects "
    "unrecognised values. work_orders.assigned_name is likewise free text rather than a user "
    "FK, so the All Inventory employee filter matches on name. These are the joints where data "
    "quality depends on validation code rather than the database."
)

INTEGRATIONS = [
    ("/partner portal", "Inbound, dealer-authenticated",
     "Trade Partner B2B — listings, bids, bookings, payment proof"),
    ("/care/api/v1", "Inbound, API key", "Post-sale device pairing, diagnostics, support tickets"),
    ("/iqc/api · /stress", "Inbound from OxyQC agent",
     "Hardware auto-detection and stress-test results from the device itself"),
    ("/whatsapp", "Bidirectional", "Dealer broadcast, quotes, market availability"),
    ("Webhooks", "Outbound", "Domain events to registered subscribers"),
    ("/bulk-upload", "Inbound, CSV", "Lots, IQC devices, telecalling, spare parts, locations"),
]

CAVEATS = (
    "Module names and grouping are accurate to the registry. Flow steps are derived from real "
    "foreign keys, so a link means a genuine referential relationship — but the absence of a "
    "link does not prove the absence of coupling. Service-level and string-based links exist "
    "that the FK graph cannot see, and the soft links noted above are the ones known to matter. "
    "Route counts reflect declared handlers, not reachable pages, since some routers serve "
    "fragments and APIs rather than screens."
)


# ── PDF ───────────────────────────────────────────────────────────────────────

def _pdf_styles():
    ss = getSampleStyleSheet()
    ss.add(ParagraphStyle("MMTitle", parent=ss["Title"], fontName="Helvetica-Bold",
                          fontSize=22, textColor=_NAVY, spaceAfter=6, alignment=0))
    ss.add(ParagraphStyle("MMSub", parent=ss["Normal"], fontName="Helvetica",
                          fontSize=11, textColor=_MUTED, spaceAfter=2))
    ss.add(ParagraphStyle("MMSection", parent=ss["Heading1"], fontName="Helvetica-Bold",
                          fontSize=13, textColor=colors.white, spaceBefore=14, spaceAfter=8,
                          backColor=_NAVY, borderPadding=(6, 6, 6, 6), leading=17))
    ss.add(ParagraphStyle("MMFlowTitle", parent=ss["Normal"], fontName="Helvetica-Bold",
                          fontSize=11, textColor=_NAVY, spaceBefore=8, spaceAfter=3))
    ss.add(ParagraphStyle("MMBody", parent=ss["Normal"], fontName="Helvetica",
                          fontSize=9.5, leading=13.5, spaceAfter=4))
    ss.add(ParagraphStyle("MMStep", parent=ss["Normal"], fontName="Helvetica",
                          fontSize=9, leading=12.5, leftIndent=10, spaceAfter=1))
    ss.add(ParagraphStyle("MMNote", parent=ss["Normal"], fontName="Helvetica-Oblique",
                          fontSize=8.5, leading=12, textColor=_MUTED, spaceAfter=4))
    ss.add(ParagraphStyle("MMCell", parent=ss["Normal"], fontName="Helvetica",
                          fontSize=8.5, leading=11.5))
    ss.add(ParagraphStyle("MMCellHead", parent=ss["Normal"], fontName="Helvetica-Bold",
                          fontSize=8.5, leading=11.5, textColor=colors.white))
    ss.add(ParagraphStyle("MMMono", parent=ss["Normal"], fontName="Courier",
                          fontSize=8, leading=11))
    ss.add(ParagraphStyle("MMFooter", parent=ss["Normal"], fontName="Helvetica-Oblique",
                          fontSize=8, textColor=_MUTED))
    return ss


def _grid_table(rows, col_widths, styles, header=True, zebra=True):
    data = []
    for i, row in enumerate(rows):
        style = styles["MMCellHead"] if (header and i == 0) else styles["MMCell"]
        data.append([Paragraph(str(c), style) for c in row])
    cmds = [
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#D5DCE6")),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]
    if header:
        cmds.append(("BACKGROUND", (0, 0), (-1, 0), _NAVY))
    if zebra:
        start = 1 if header else 0
        for r in range(start, len(data)):
            if (r - start) % 2 == 1:
                cmds.append(("BACKGROUND", (0, r), (-1, r), _LIGHT))
    return Table(data, colWidths=col_widths, style=TableStyle(cmds), repeatRows=1 if header else 0)


def build_module_map_pdf(*, generated_on: str, source_ref: str = SOURCE_REF) -> bytes:
    """Renders the Module & Interconnectivity Map as a styled PDF and returns the bytes."""
    buf = BytesIO()
    width = 6.9 * inch
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER,
        leftMargin=0.75 * inch, rightMargin=0.75 * inch,
        topMargin=0.75 * inch, bottomMargin=0.75 * inch,
        title="OxyPC Inventory — Module & Interconnectivity Map",
    )
    ss = _pdf_styles()
    story = []

    # Cover block
    story.append(Paragraph("OxyPC Inventory — Module &amp; Interconnectivity Map", ss["MMTitle"]))
    story.append(Paragraph(
        "Every module the application exposes, grouped by the part of the business it serves — "
        "and the flows that carry a device, a part, or an order between them.", ss["MMSub"]))
    story.append(Paragraph(
        f"Document code: MIM-1.0 &nbsp;|&nbsp; Source: {source_ref} &nbsp;|&nbsp; "
        f"Generated: {generated_on}", ss["MMSub"]))
    story.append(Spacer(1, 4))
    story.append(Table([[""]], colWidths=[width], rowHeights=[2],
                       style=TableStyle([("BACKGROUND", (0, 0), (-1, -1), _ACCENT)])))
    story.append(Spacer(1, 10))

    # Headline stats strip
    stat_cells = [[Paragraph(f"<font size=13 color='#1F3864'><b>{v}</b></font><br/>"
                             f"<font size=7.5 color='#5A6472'>{k.upper()}</font>", ss["MMCell"])
                   for v, k in HEADLINE_STATS]]
    story.append(Table(stat_cells, colWidths=[width / 6.0] * 6, style=TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), _LIGHT),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 7),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
        ("BOX", (0, 0), (-1, -1), 0.4, colors.HexColor("#D5DCE6")),
        ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.white),
    ])))
    story.append(Spacer(1, 6))
    story.append(Paragraph(
        "Module names come from PERM_MODULES in routers/master.py — the same registry that "
        "drives the permission matrix, sidebar labels and page titles. Flows come from the "
        "foreign-key graph across all model files plus the stage machine in "
        "services/control_engine.py.", ss["MMNote"]))

    # 1. Spine
    story.append(Paragraph("1. The spine — device lifecycle", ss["MMSection"]))
    story.append(Paragraph(SPINE_INTRO, ss["MMBody"]))
    for step in SPINE_STAGES:
        story.append(Paragraph(f"•&nbsp; {step}", ss["MMStep"]))
    story.append(Spacer(1, 4))
    story.append(Paragraph(SPINE_NOTE, ss["MMNote"]))

    # 2. Module registry
    story.append(Paragraph("2. Module registry — all 88", ss["MMSection"]))
    story.append(Paragraph(
        "Grouped by the business function they serve. Group membership is editorial — the "
        "registry itself is a flat list; grouping is how the sidebar is arranged and how "
        "permissions tend to be granted in practice.", ss["MMBody"]))
    rows = [["Group", "#", "Modules"]]
    for group, mods in MODULE_GROUPS:
        rows.append([group, str(len(mods)), "&nbsp;· ".join(mods)])
    story.append(_grid_table(rows, [1.35 * inch, 0.3 * inch, width - 1.65 * inch], ss))

    # 3. Flows
    story.append(Paragraph("3. Interconnectivity flows — eight cross-module paths", ss["MMSection"]))
    story.append(Paragraph(
        "Numbered because they are genuine sequences — each hand-off has a direction, and later "
        "flows depend on state the earlier ones create.", ss["MMBody"]))
    for code, title, summary, steps, joints, warning in FLOWS:
        block = [Paragraph(f"{code} &nbsp;{title}", ss["MMFlowTitle"]),
                 Paragraph(summary, ss["MMBody"])]
        for step in steps:
            block.append(Paragraph(f"•&nbsp; {step}", ss["MMStep"]))
        block.append(Spacer(1, 3))
        block.append(Paragraph(f"<b>Key joints:</b> {joints}", ss["MMNote"]))
        if warning:
            block.append(Table(
                [[Paragraph(warning, ss["MMCell"])]], colWidths=[width],
                style=TableStyle([
                    ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F7E6E3")),
                    ("BOX", (0, 0), (-1, -1), 0.4, _RED),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ])))
        block.append(Spacer(1, 6))
        story.append(KeepTogether(block))

    # 4. Control plane
    story.append(Paragraph("4. Cross-cutting control plane", ss["MMSection"]))
    story.append(Paragraph(
        "These are not modules and have no sidebar entry, but nearly every write in the system "
        "goes through at least one of them.", ss["MMBody"]))
    rows = [["Service", "Responsibility", "Reaches"]]
    rows.extend([[s, r, x] for s, r, x in SERVICES])
    story.append(_grid_table(rows, [1.5 * inch, width - 3.0 * inch, 1.5 * inch], ss))
    story.append(Spacer(1, 4))
    story.append(Paragraph(GOVERNANCE_NOTE, ss["MMNote"]))

    # 5. Coupling
    story.append(Paragraph("5. Coupling — where change is expensive", ss["MMSection"]))
    story.append(Paragraph(
        "A high inbound foreign-key count means a schema change there ripples widely. These six "
        "tables carry 76 of the 142 FK edges between them.", ss["MMBody"]))
    rows = [["Table", "Referenced by", "What that means in practice"]]
    rows.extend([[t, n, m] for t, n, m in COUPLING])
    story.append(_grid_table(rows, [1.2 * inch, 0.9 * inch, width - 2.1 * inch], ss))
    story.append(Spacer(1, 4))
    story.append(Paragraph(f"<b>Soft links — coupling without a foreign key:</b> {SOFT_LINKS}",
                           ss["MMNote"]))

    # 6. Integration surface
    story.append(Paragraph("6. Integration surface", ss["MMSection"]))
    rows = [["Surface", "Direction", "Purpose"]]
    rows.extend([[s, d, p] for s, d, p in INTEGRATIONS])
    story.append(_grid_table(rows, [1.35 * inch, 1.5 * inch, width - 2.85 * inch], ss))

    # 7. Caveats
    story.append(Paragraph("7. Caveats — what this map is and is not", ss["MMSection"]))
    story.append(Paragraph(CAVEATS, ss["MMBody"]))

    story.append(Spacer(1, 10))
    story.append(Table([[""]], colWidths=[width], rowHeights=[0.75],
                       style=TableStyle([("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#CCCCCC"))])))
    story.append(Spacer(1, 4))
    story.append(Paragraph(
        f"{source_ref} · {GENERATED_LABEL} · generated {generated_on}. Downloaded on demand "
        "from the OxyPC Inventory QA Dashboard — regenerate after any major architectural change.",
        ss["MMFooter"]))

    doc.build(story)
    return buf.getvalue()


# ── PPTX ──────────────────────────────────────────────────────────────────────

def build_module_map_pptx(*, generated_on: str, source_ref: str = SOURCE_REF) -> bytes:
    """Renders the Module & Interconnectivity Map as a 16:9 deck and returns the bytes."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    navy = RGBColor(0x1F, 0x38, 0x64)
    accent = RGBColor(0x2E, 0x75, 0xB6)
    light = RGBColor(0xEE, 0xF3, 0xF8)
    muted = RGBColor(0x5A, 0x64, 0x72)
    red = RGBColor(0x9E, 0x3A, 0x2C)
    white = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    def textbox(slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        return tf

    def para(tf, text, *, size=14, bold=False, color=None, space_after=4, first=False,
             align=PP_ALIGN.LEFT, font="Calibri"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = text
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color if color is not None else RGBColor(0x22, 0x22, 0x22)
        return p

    def band(slide, title):
        """Navy header band across the top of a content slide."""
        from pptx.enum.shapes import MSO_SHAPE
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.95))
        bar.fill.solid()
        bar.fill.fore_color.rgb = navy
        bar.line.fill.background()
        bar.shadow.inherit = False
        tf = bar.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.6)
        para(tf, title, size=24, bold=True, color=white, first=True)

    def add_table(slide, rows, left, top, width, col_ratios, *, font_size=10.5,
                  header_size=10.5, row_height=Inches(0.32)):
        n_rows, n_cols = len(rows), len(rows[0])
        shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows)
        table = shape.table
        total = float(sum(col_ratios))
        for i, ratio in enumerate(col_ratios):
            table.columns[i].width = Emu(int(width * (ratio / total)))
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(val)
                cell.margin_left = Inches(0.07)
                cell.margin_right = Inches(0.07)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                cell.fill.solid()
                if r == 0:
                    cell.fill.fore_color.rgb = navy
                else:
                    cell.fill.fore_color.rgb = light if r % 2 == 0 else white
                p = cell.text_frame.paragraphs[0]
                p.space_after = Pt(0)
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(header_size if r == 0 else font_size)
                run.font.bold = (r == 0)
                run.font.name = "Calibri"
                run.font.color.rgb = white if r == 0 else RGBColor(0x22, 0x22, 0x22)
        return table

    # — Slide 1: cover
    s = prs.slides.add_slide(blank)
    from pptx.enum.shapes import MSO_SHAPE
    hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(4.1))
    hero.fill.solid()
    hero.fill.fore_color.rgb = navy
    hero.line.fill.background()
    hero.shadow.inherit = False
    tf = hero.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.9)
    tf.margin_top = Inches(1.0)
    para(tf, "OxyPC Inventory", size=20, color=RGBColor(0x9E, 0xC4, 0xE8), first=True)
    para(tf, "Module & Interconnectivity Map", size=40, bold=True, color=white, space_after=10)
    para(tf, "Every module the application exposes — and the flows that carry a device, "
             "a part, or an order between them.", size=15, color=RGBColor(0xC9, 0xD8, 0xE8))
    para(tf, f"Document code MIM-1.0  ·  Source {source_ref}  ·  Generated {generated_on}",
         size=11, color=RGBColor(0x9E, 0xC4, 0xE8))

    # stat strip
    strip_top = Inches(4.45)
    box_w = Emu(int((SW - Inches(1.8)) / 6))
    for i, (val, label) in enumerate(HEADLINE_STATS):
        bx = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(Inches(0.9) + i * box_w)),
                                strip_top, Emu(int(box_w * 0.94)), Inches(1.0))
        bx.fill.solid()
        bx.fill.fore_color.rgb = light
        bx.line.color.rgb = RGBColor(0xD5, 0xDC, 0xE6)
        bx.shadow.inherit = False
        t = bx.text_frame
        t.word_wrap = True
        para(t, val, size=22, bold=True, color=navy, align=PP_ALIGN.CENTER, first=True, space_after=0)
        para(t, label.upper(), size=9, color=muted, align=PP_ALIGN.CENTER)

    tf = textbox(s, Inches(0.9), Inches(5.8), SW - Inches(1.8), Inches(1.2))
    para(tf, "Module names come from PERM_MODULES in routers/master.py — the same registry that "
             "drives the permission matrix, sidebar labels and page titles. Flows come from the "
             "foreign-key graph across all model files plus the stage machine in "
             "services/control_engine.py.", size=11, color=muted, first=True)

    # — Slide 2: the spine
    s = prs.slides.add_slide(blank)
    band(s, "The spine — device lifecycle")
    tf = textbox(s, Inches(0.6), Inches(1.2), SW - Inches(1.2), Inches(1.0))
    para(tf, SPINE_INTRO, size=12.5, color=muted, first=True)
    tf = textbox(s, Inches(0.6), Inches(2.35), SW - Inches(1.2), Inches(3.0))
    for i, step in enumerate(SPINE_STAGES):
        para(tf, "▸  " + step, size=14, first=(i == 0), space_after=8)
    note = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.85),
                              SW - Inches(1.2), Inches(1.05))
    note.fill.solid()
    note.fill.fore_color.rgb = light
    note.line.color.rgb = RGBColor(0xD5, 0xDC, 0xE6)
    note.shadow.inherit = False
    t = note.text_frame
    t.word_wrap = True
    t.margin_left = Inches(0.2)
    para(t, SPINE_NOTE, size=11, color=muted, first=True)

    # — Slides 3–4: module registry (split so the tables stay readable)
    for idx, chunk in enumerate((MODULE_GROUPS[:6], MODULE_GROUPS[6:])):
        s = prs.slides.add_slide(blank)
        band(s, f"Module registry — all 88 ({idx + 1} of 2)")
        rows = [["Group", "#", "Modules"]]
        for group, mods in chunk:
            rows.append([group, str(len(mods)), " · ".join(mods)])
        add_table(s, rows, Inches(0.6), Inches(1.3), SW - Inches(1.2), [2.4, 0.5, 9.0],
                  font_size=10, row_height=Inches(0.55))

    # — Slides 5–12: one per flow
    for code, title, summary, steps, joints, warning in FLOWS:
        s = prs.slides.add_slide(blank)
        band(s, f"{code} — {title}")
        tf = textbox(s, Inches(0.6), Inches(1.15), SW - Inches(1.2), Inches(0.9))
        para(tf, summary, size=12.5, color=muted, first=True)
        tf = textbox(s, Inches(0.6), Inches(2.15), Inches(7.4), Inches(3.6))
        for i, step in enumerate(steps):
            para(tf, "▸  " + step, size=12.5, first=(i == 0), space_after=7)
        side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(2.15),
                                  Inches(4.45), Inches(3.6))
        side.fill.solid()
        side.fill.fore_color.rgb = light
        side.line.color.rgb = RGBColor(0xD5, 0xDC, 0xE6)
        side.shadow.inherit = False
        t = side.text_frame
        t.word_wrap = True
        t.margin_left = Inches(0.18)
        t.margin_top = Inches(0.14)
        para(t, "KEY JOINTS", size=10, bold=True, color=accent, first=True)
        para(t, joints, size=11, color=RGBColor(0x33, 0x33, 0x33))
        if warning:
            wbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.95),
                                      SW - Inches(1.2), Inches(1.0))
            wbox.fill.solid()
            wbox.fill.fore_color.rgb = RGBColor(0xF7, 0xE6, 0xE3)
            wbox.line.color.rgb = red
            wbox.shadow.inherit = False
            t = wbox.text_frame
            t.word_wrap = True
            t.margin_left = Inches(0.2)
            para(t, warning, size=11, color=RGBColor(0x6B, 0x27, 0x1D), first=True)

    # — Control plane
    s = prs.slides.add_slide(blank)
    band(s, "Cross-cutting control plane")
    rows = [["Service", "Responsibility", "Reaches"]]
    rows.extend([[a, b, c] for a, b, c in SERVICES])
    add_table(s, rows, Inches(0.6), Inches(1.15), SW - Inches(1.2), [2.6, 6.2, 2.6],
              font_size=10, row_height=Inches(0.36))
    tf = textbox(s, Inches(0.6), Inches(6.5), SW - Inches(1.2), Inches(0.8))
    para(tf, GOVERNANCE_NOTE, size=10.5, color=muted, first=True)

    # — Coupling
    s = prs.slides.add_slide(blank)
    band(s, "Coupling — where change is expensive")
    rows = [["Table", "Referenced by", "What that means in practice"]]
    rows.extend([[a, b, c] for a, b, c in COUPLING])
    add_table(s, rows, Inches(0.6), Inches(1.25), SW - Inches(1.2), [2.0, 1.5, 8.4],
              font_size=11, row_height=Inches(0.45))
    tf = textbox(s, Inches(0.6), Inches(4.7), SW - Inches(1.2), Inches(2.0))
    para(tf, "SOFT LINKS — COUPLING WITHOUT A FOREIGN KEY", size=10, bold=True,
         color=accent, first=True)
    para(tf, SOFT_LINKS, size=11.5, color=RGBColor(0x33, 0x33, 0x33))

    # — Integration surface
    s = prs.slides.add_slide(blank)
    band(s, "Integration surface")
    rows = [["Surface", "Direction", "Purpose"]]
    rows.extend([[a, b, c] for a, b, c in INTEGRATIONS])
    add_table(s, rows, Inches(0.6), Inches(1.35), SW - Inches(1.2), [2.4, 2.8, 6.7],
              font_size=11, row_height=Inches(0.5))

    # — Caveats / close
    s = prs.slides.add_slide(blank)
    band(s, "Caveats — what this map is and is not")
    tf = textbox(s, Inches(0.6), Inches(1.5), SW - Inches(1.2), Inches(3.0))
    para(tf, CAVEATS, size=14, first=True)
    tf = textbox(s, Inches(0.6), Inches(6.2), SW - Inches(1.2), Inches(0.8))
    para(tf, f"{source_ref} · {GENERATED_LABEL} · generated {generated_on}",
         size=10.5, color=muted, first=True)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
